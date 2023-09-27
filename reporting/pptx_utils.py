import logging
logger = logging.getLogger(__name__)
import datetime

# presentation manipulation
import pptx
import copy

import reporting.plot_utils as plot
import module.data_transformation as transform
# text manipulation
import re
from num2words import num2words

# image manipulation
from io import BytesIO
import pandas as pd

# pptx plot
from pptx.chart.data import CategoryChartData	# Classes providing reference data types


PLACEHOLDER_PATTERNS: dict = {
    "default": "\[ph_([a-zA-Z0-9_]+)\]"
}

def solution_description_slide(slide: pptx.slide.Slide, rows: dict) -> None:
    for idx, row_data in rows.items():
        placeholder_pattern = f'\[row{idx}_([a-zA-Z0-9_]+)\]'
        fill_text_placeholders(slide, row_data, placeholder_pattern)

def qualitative_plots_slide(presentation: pptx.Presentation, template_slide: pptx.slide.Slide, num_page_add: int, img: dict) -> None:
    #slides: list[pptx.slide.Slide] = []
    for idx in range(num_page_add):
        slide: pptx.slide.Slide = create_empty_slides(presentation, template_slide, pages_to_add=1)[0]
        #slides.append(slide)
        fill_single_image_placeholders(slide,list(img.items())[idx][1] )
    
 

def fin_indicator_slide(template_slide: pptx.slide.Slide,df_fin_performance: pd.DataFrame) -> None:
    # transform the data into chart_data by group of fin indicator
    chart_data_list: list = transform.fin_indicator_plot_data(df_fin_performance)
    # insert the charts to the pptx
    pptx_charts(template_slide,chart_data_list)
    
def competitor_slide(template_slide: pptx.slide.Slide,competitor_data: pd.DataFrame, df_fin_performance: pd.DataFrame) -> None:
    
    competitor_data, competitor_name = transform.competitor_data_to_ch(df_fin_performance, competitor_data)
    chart_data_list: list = transform.fin_competitor_plot_data(competitor_data, competitor_name)
    pptx_charts(template_slide, chart_data_list)
    
def interviewee_gap_slide(slide: pptx.slide.Slide, pic_rows: pd.DataFrame, text_rows: dict)-> None:
    
    for idx, row_data in text_rows.items():
        placeholder_pattern = f'\[row{idx}_([a-zA-Z0-9_]+)\]'
        fill_text_placeholders(slide, row_data, placeholder_pattern)
    
    
    img_buffers: dict = {} 
    # loop over unique values in column "modules"
    for value in pic_rows['module'].unique():
        # create new DataFrame for this value of "A"
        df = pic_rows[pic_rows['module'] == value].reset_index()
        img_buffers[value] = plot.interviewee_plots(df) 
    
    for key in  img_buffers:
        fill_single_image_placeholders(slide,img_buffers[key] )

# define function to sum values in dictionary column
def sum_dict_values(dicts):
    return sum(d.values() for d in dicts)
    

def qualitative_questions_detail_slide(slide: pptx.slide.Slide, image_buffer: BytesIO) -> None:
    fill_single_image_placeholders(slide, image_buffer)


def fin_indicator_sensitivity_slide(slide: pptx.slide.Slide, image_buffer: BytesIO) -> None:
    fill_single_image_placeholders(slide, image_buffer)


def solution_priority_matrix_slide(slide: pptx.slide.Slide, image_buffer: BytesIO) -> None:
    fill_single_image_placeholders(slide, image_buffer)


def solution_roadmap_slide(slide: pptx.slide.Slide, rows: dict) -> None:
    for idx, row_data in rows.items():
        placeholder_pattern = f'\[row{idx}_([a-zA-Z0-9_]+)\]'
        fill_text_placeholders(slide, row_data, placeholder_pattern)


def solution_roi_slide(slide: pptx.slide.Slide, rows: dict) -> None:
    # data manipulation
    # execute
    for idx, row_data in rows.items():
        placeholder_pattern = f'\[row{idx}_([a-zA-Z0-9_]+)\]'
        fill_text_placeholders(slide, row_data, placeholder_pattern)


def qualitative_gap_slide(slide: pptx.slide.Slide, rows: dict) -> None:
    # data manipulation
    # execute
    for idx, row_data in rows.items():
        placeholder_pattern = f'\[row{idx}_([a-zA-Z0-9_]+)\]'
        fill_text_placeholders(slide, row_data, placeholder_pattern)


def company_slide(slide: pptx.slide.Slide, company: dict) -> None:
    # data manipulation
    # 資本額: 三億, 三億-五億
    company['capital_text'] = '-'.join([
        num2words(str(x), lang='ja') 
        for x in (company['capital_min'], company['capital_max']) if x is not None
        ])

    # 員工人數: 26人, 26-52人
    company['employee_text'] = '-'.join([
        f'{str(int(x))} 人'
        for x in (company['employee_min'], company['employee_max']) if x is not None
        ])

    # execute
    fill_text_placeholders(slide, data=company, placeholder_pattern=PLACEHOLDER_PATTERNS['default'])


def strategy_slide(slide: pptx.slide.Slide, strategy: dict) -> None:
    # data manipulation
    # execute
    fill_text_placeholders(slide, data=strategy, placeholder_pattern=PLACEHOLDER_PATTERNS['default'])


def industry_slides( presentation: pptx.Presentation, template_slide: pptx.slide.Slide, 
                                cases_per_slide: int, industries: dict, cases: pd.DataFrame.groupby)  -> None:

    for industry_id, industry_data in industries.items():
        # calculate how many pages need to be add.
        case_count = industry_data['case_count']
        num = case_count // cases_per_slide
        industry_data['pages_to_add'] = num + 1 if case_count % cases_per_slide != 0 else num
        
        # create template slides.
        slides: list[pptx.slide.Slide] = create_empty_slides(presentation, template_slide, pages_to_add=industry_data['pages_to_add'])
        industry_data['slides'] = slides
    
    # for each industry, fill case into slide
    for industry_id, industry_data in industries.items():
        # variable init.
        industry_cases: list[dict] = cases.get_group(industry_id).to_dict(orient='records')
        slides: list[pptx.slide.Slide] = industry_data['slides']
        is_all_cases_filled = False
        
        for empty_slide in slides:
            # fill industry data into slide.
            logger.debug(f'now doing slide: {empty_slide}')
            fill_text_placeholders(empty_slide, data=industry_data, placeholder_pattern=PLACEHOLDER_PATTERNS['default'])

            # for each case slot in slide.
            slot_id = 0
            while slot_id < cases_per_slide and not is_all_cases_filled:
                logger.debug(f'now doing slot: {slot_id}')
                case_data = industry_cases.pop(0)
                fill_case(empty_slide, case_data, placeholder_pattern=f'\[case{slot_id}_ph_([a-zA-Z0-9_]+)\]')

                if len(industry_cases) == 0:
                    is_all_cases_filled = True
                
                slot_id += 1
            
            # if no more case to fill, break.
            if is_all_cases_filled:
                break


def create_empty_slides(presentation: pptx.Presentation, source_slide: pptx.slide.Slide, pages_to_add: int) -> list[pptx.slide.Slide]:
    # add 客戶產業數位轉型重點與建議 slides to presentation, and return list of slide objects.
    slides = []
    for i in range(pages_to_add):
        slides.append(duplicate_slide(presentation, source_slide))
    return slides


def fill_case(industry_slide: pptx.slide.Slide, case_data: dict, placeholder_pattern: str) -> None:

    case_image_inserted = False

    for shape in industry_slide.shapes:

        # 只針對 placeholder 進行操作
        if not shape.is_placeholder:
            continue

        # case image: picture placeholder 操作
        if isinstance(shape, pptx.shapes.placeholder.PicturePlaceholder):
            
            if case_image_inserted:                                         # only insert case pitcure once.
                continue

            image_buffer = BytesIO(case_data['case_img_blob'])              # read bytes into file like object: BytesIO.
            shape.insert_picture(image_buffer)                              # after this, picture placeholder would become invalid, and shape object will become pptx.shapes.placeholder.PlaceholderPicture
            case_image_inserted = True
            continue

        # case text: text placeholder 操作
        if not shape.has_text_frame:
            continue
        
        for paragraph in shape.text_frame.paragraphs:
            
            # 若沒有欲取代字串，繼續看下一個 paragraph.             
            placeholder_match = re.search(placeholder_pattern, paragraph.text) # [case0_ph_case_text]
            if placeholder_match is None:
                continue

            # 取代字串
            placeholder_text = placeholder_match.group(0)                   # [case0_ph_case_text]
            column_id = placeholder_match.group(1)                          # case_text
            new_text = paragraph.text.replace(placeholder_text, case_data[column_id])
            replace_paragraph_text(paragraph, new_text)


def fill_single_image_placeholders(slide: pptx.slide.Slide, image_buffer: BytesIO) -> None:
    for shape in slide.shapes:

        if not isinstance(shape, pptx.shapes.placeholder.PicturePlaceholder):
            continue

        # picture placeholder 操作
        shape.insert_picture(image_buffer)                              # after this, picture placeholder would become invalid, and shape object will become pptx.shapes.placeholder.PlaceholderPicture
        return



def fill_text_placeholders(slide: pptx.slide.Slide, data: dict, placeholder_pattern: str) -> None:
    # https://magenta-fern-2ff.notion.site/Placeholder-b2b98d29d2554c89a357593ec0e9e153
    for shape in slide.shapes:

        # 尋找 placeholder 圖形       
        if not shape.is_placeholder:
            continue

        # text placeholder 操作
        if not shape.has_text_frame:
            continue
        
        for paragraph in shape.text_frame.paragraphs:
            
            # 若沒有欲取代字串，繼續看下一個 paragraphs.
            placeholder_match = re.search(placeholder_pattern, paragraph.text)
            if placeholder_match is None:
                continue

            # 取代字串
            placeholder_text = placeholder_match.group(0) # [ph_company_text]
            column_id = placeholder_match.group(1)        # company_text
            new_text = paragraph.text.replace(placeholder_text, str(data[column_id]))
            replace_paragraph_text(paragraph, new_text)


def replace_paragraph_text(paragraph, new_text) -> None:
    # replace text while preserve formatting: https://github.com/scanny/python-pptx/issues/285
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)

    paragraph.runs[0].text = new_text


def duplicate_slide(presentation: pptx.Presentation, source: pptx.slide.Slide) -> pptx.slide.Slide:
    # https://github.com/scanny/python-pptx/issues/132
    destination = presentation.slides.add_slide(source.slide_layout)

    for shape in destination.shapes:
        shape.element.getparent().remove(shape.element)
    
    for shape in source.shapes:
        new_shape = copy.deepcopy(shape.element)
        destination.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
    
    for relationship in source.part.rels:
        target = relationship._target

        if "noteSlide" in relationship.reltype: continue

        destination.part.rels.get_or_add(relationship.reltype, relationship._target)

    return destination

def  cover_slide(slide: pptx.slide.Slide, company: dict) -> None:
    # 獲取當下月份年份
    now = datetime.datetime.now()
    formatted_time: str = now.strftime(("%m-%Y"))
    company["MM_YYYY"] = formatted_time
    
    # execute
    fill_text_placeholders(slide, data=company, placeholder_pattern=PLACEHOLDER_PATTERNS['default'])
    
def pptx_charts(template_slide: pptx.slide.Slide, chart_data_list: list) -> None:
    slide = template_slide
    
    # Loop through chart data and insert charts
    for i, chart_data in enumerate(chart_data_list):
        # Select the chart placeholder on the slide
        chart = None
        
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            
            for series in shape.chart.series:
                if series.name != chart_data["df_name_ch"]:
                    continue
                chart = shape.chart
                break
                
        # Access the chart data and modify it
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = chart_data["categories"]
        # Data setting
        if "(%)" in chart_data["df_name_ch"]:
            chart_data["series"] = tuple( i*100  for i in chart_data["series"]) 
        elif "(千元)" in chart_data["df_name_ch"]:
            chart_data["series"] = tuple( i/1000 for i in chart_data["series"]) 
        elif "(天)" in chart_data["df_name_ch"]:
            chart_data["series"] = tuple( round(i, 0) for i in chart_data["series"])  
        chart_data_obj.add_series(chart_data["df_name_ch"], chart_data["series"])
 

        # Update the chart object to reflect the changes
        chart.replace_data(chart_data_obj)

