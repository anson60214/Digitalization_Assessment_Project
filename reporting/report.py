
import sqlalchemy as database
import pandas as pd
import json

# presentation manipulation
import pptx
import reporting.pptx_utils as report
import reporting.plot_utils as plot

# database model
from db.model_stg import *
import db.repository_stg as repo
import module.data_transformation as transform

from io import BytesIO

from PIL import Image


TEST_DATA: dict = {
    "company_industry_l_id": "F"
}

TEST_PRESENTATION_TEMPLATE_NAME: str = "reporting\\數位轉型服務_Final Report_Template_0303_v1.1.pptx"

TEMPLATE_SLIDE_MAP: dict = {
    "數位轉型專案封面": 0,
    "公司基本資料": 5,
    "客戶主要的三大發展策略": 10,
    #"客戶產業數位轉型重點與建議": 5,
    "客戶八大財務角度分析": 11,
    "與其他競爭者相比的八大財務分析": 12,
    "質化問卷分數落差分析": 13,
    "PwC針對四大面向提供相應的建議-1": 14,
    "PwC針對四大面向提供相應的建議-2": 15,
    #"財務指標表現": 8,
    #"財務敏感度影響分析": 9,
    #"解決方案優先順序矩陣圖": 11,
    "PwC潛在建議方案": 19,
    "請客戶就潛在方案進行排序": 20,
    "解決方案 ROI 排名": 21,
    "解決方案規劃建議時程": 22,
    "受訪者差異分析: 數位人才": 27,
    "受訪者差異分析: 顧客體驗": 28,
    "受訪者差異分析: 數位營運": 29,
    "受訪者差異分析: 數位科技": 30,
}

CASES_PER_SLIDE = 2
INDICATOR_PER_SLIDE = 4

# 0. 封面資訊
# 1. 公司基本資料 
# 2. 客戶主要的三大發展策略
# 3. 客戶產業數位轉型重點與建議
# 4. 質化問卷分數落差分析
# 5. PwC針對四大面向提供相應的建議
# 6. 與其他競爭者相比的八大財務分析
# 6. 財務敏感度影響分析, plot
# 7. PwC潛在建議方案, 請客戶就潛在方案進行排序, plot -> 解決方案優先順序矩陣圖
# 8. 解決方案 ROI 排名
# 9. 解決方案規劃建議時程
# 11. plot: 質化明細
# 12. 受訪者差異分析
def generate_report(conn: database.engine, input_tables: dict[str, pd.DataFrame], calculated_tables: dict[str, pd.DataFrame]) -> BytesIO:

    # df_solution: 用於 7, 8, 9, 10，綜合分數前八名解決方案，包含 ROI
    # df_year_data: 用於 5, 6製圖，包含各財務指標3年資料.
    # df_qualitative_result: 用於 4, 11，使用者問卷作答明細，包含aspect, module, question, actual(as-is), target(to-be), gap, ql_score
    # df_trend: 用於 5，包含 sensitivity analysis, performance_gap_impact_cashflow.

    df_company_data = input_tables['df_company_data']
    df_form_data = input_tables['df_form_data']
    strategy_id = transform.get_form_value(df_company_data, "STRAT").split(".")[0]
    company_id = transform.get_form_value(df_company_data, "company_id")

    df_competitor = input_tables['df_competitor']

    df_solution = calculated_tables['解決方案前十名與ROI']
    df_year_data = calculated_tables['三年財務指標運算結果']
    df_qualitative_result = calculated_tables['質化分析運算結果']
    df_trend = calculated_tables['趨勢現金流與敏感度分析']
    
    # ------------------------------------------------------------------------------------------------------------------
    
    # 1. 公司基本資料
    target_company: dict = repo.get_company_data(conn, company_id)[0]
    
    # 2. 客戶主要的三大發展策略
    target_strategy: dict = repo.get_strategy_weight(conn, strategy_id)[0]

    # 補上company name in dict
    target_strategy["company_text"] = target_company["company_text"]
    
    # 3. 客戶產業數位轉型重點與建議
    #industries_cases: pd.DataFrame = repo.get_industries_cases(conn, TEST_DATA['company_industry_l_id'])
    #industries: dict = transform.industry_with_case_count(industries_cases)
    #cases = industries_cases.groupby('industry_id')
    
    # 4. 質化問卷分數落差分析
    # merge question into df_qualitative_result
    qualitative_question: pd.DataFrame = repo.get_dim_qualitative_question(conn)[["question_id","question"]]
    df_qualitative_result_question =  pd.merge(df_qualitative_result, qualitative_question, on='question_id', how='left')
    qualitative_top10_gap: dict = transform.qualitative_top10_gap(df_qualitative_result_question)

    # 5. PwC針對四大面向提供相應的建議
    qualitative_top3_gap_aspect1 : dict = transform.qualitative_top3_gap(df_qualitative_result_question, "數位人才", "顧客體驗")
    qualitative_top3_gap_aspect2 : dict = transform.qualitative_top3_gap(df_qualitative_result_question, "數位營運", "新科技")
    
    # 6. 財務指標表現, plot
    #   作圖呈現三年財務指標數據變化，使用 df_fin_performance 作為作圖數據。
    df_dim_quantative_index = repo.get_dim_quantative_index(conn)
    df_fin_performance, df_fin_sensitivity = transform.fin_indicator_calculation_result(df_year_data, df_dim_quantative_index)
    #img_buffers = { 
        #indicator_name: plot.fin_performance(indicator_name, df) 
        #for indicator_name, df in df_fin_performance.groupby('fin_indicator_text_en')}
    #df_trend['png_buffer'] = df_trend['fin_indicator_text_en'].map(img_buffers)
    #fin_indicator_data: list[dict] = transform.fin_indicator_data(df_trend)
    
    # 競爭對手財務指標
    competitor_data = transform.competitor_data(df_competitor, df_year_data, target_company["company_text"])


    # 6. 財務敏感度影響分析
    fin_sensitivity_plot_png: BytesIO = plot.fin_sensitivity(df_fin_sensitivity)
    
    # 7. 解決方案優先順序矩陣圖
    solution_priority_matrix_png: BytesIO = plot.solution_priority_matrix(df_solution)

    # 8. Solution Roi, 9. Solution Roadmap, 10. Solution Description
    solution_ranking: dict = transform.solution_ranking(df_solution)

    # 11. plot: 質化明細
    # qualitative_plot_png: BytesIO = plot.qualitative_detail(df_qualitative_result)
    df_qualitative_plot: pd.DataFrame = df_qualitative_result.groupby(['aspect', 'module']).mean(numeric_only=True)
    aspect_count: int = df_qualitative_plot.index.get_level_values(0).nunique()
    aspects: list = df_qualitative_plot.index.get_level_values(0).unique()
    img2_buffers:dict = { 
        aspect: plot.qualitative_detail(df_qualitative_plot, aspect, aspect_idx) 
        for aspect,aspect_idx in zip(aspects, range(aspect_count) )}
    
    # 12. 受訪者差異分析
    print(df_qualitative_result)
    df_interviewee: pd.DataFrame = pd.merge(df_form_data,  df_qualitative_result[[ "question_id", "aspect",  "module" ]], on='question_id', how='left' )
    df_interviewee: pd.DataFrame =  pd.merge(df_interviewee, qualitative_question, on='question_id', how='left')
    aspect1_qualitative_top5_gap, aspect1_df_text_diff =  transform.top_5_module_by_interviewee(df_qualitative_result_question, df_interviewee, "數位人才")
    aspect2_qualitative_top5_gap, aspect2_df_text_diff =  transform.top_5_module_by_interviewee(df_qualitative_result_question, df_interviewee, "顧客體驗")
    aspect3_qualitative_top5_gap, aspect3_df_text_diff =  transform.top_5_module_by_interviewee(df_qualitative_result_question, df_interviewee, "數位營運")
    aspect4_qualitative_top5_gap, aspect4_df_text_diff =  transform.top_5_module_by_interviewee(df_qualitative_result_question, df_interviewee, "新科技")
    
    # ------------------------------------------------------------------------------------------------------------------
    # slide generation
    presentation = pptx.Presentation(TEST_PRESENTATION_TEMPLATE_NAME)
    #pptx_template= repo.get_dim_report_template(conn)
    #presentation = pptx.Presentation(BytesIO(pptx_template.iloc[0,1]))
    
    # generate map: slide_name - slide_object
    template_slides: dict[str, pptx.slide.Slide] = {}
    for slide_name, slide_id in TEMPLATE_SLIDE_MAP.items():
        template_slides[slide_name] = presentation.slides[slide_id]
     
        
    report.cover_slide(template_slides['數位轉型專案封面'], company=target_company)
    report.company_slide(template_slides['公司基本資料'], company=target_company)
    report.strategy_slide(template_slides['客戶主要的三大發展策略'], strategy=target_strategy)
    report.fin_indicator_slide(template_slides['客戶八大財務角度分析'], df_fin_performance)
    report.competitor_slide(template_slides['與其他競爭者相比的八大財務分析'], competitor_data, df_fin_performance)
    report.qualitative_gap_slide(template_slides['質化問卷分數落差分析'], rows=qualitative_top10_gap)
    report.qualitative_gap_slide(template_slides['PwC針對四大面向提供相應的建議-1'], rows=qualitative_top3_gap_aspect1)
    report.qualitative_gap_slide(template_slides['PwC針對四大面向提供相應的建議-2'], rows=qualitative_top3_gap_aspect2)
    report.solution_description_slide(template_slides['PwC潛在建議方案'], rows=solution_ranking)
    report.solution_priority_matrix_slide(template_slides['PwC潛在建議方案'], solution_priority_matrix_png)
    report.solution_description_slide(template_slides['請客戶就潛在方案進行排序'], rows=solution_ranking)
    report.solution_roi_slide(template_slides['解決方案 ROI 排名'], rows=solution_ranking)
    report.solution_roadmap_slide(template_slides['解決方案規劃建議時程'], rows=solution_ranking)
    report.interviewee_gap_slide(template_slides['受訪者差異分析: 數位人才'], pic_rows=aspect1_qualitative_top5_gap, text_rows = aspect1_df_text_diff)
    report.interviewee_gap_slide(template_slides['受訪者差異分析: 顧客體驗'], pic_rows=aspect2_qualitative_top5_gap, text_rows = aspect2_df_text_diff)
    report.interviewee_gap_slide(template_slides['受訪者差異分析: 數位營運'], pic_rows=aspect3_qualitative_top5_gap, text_rows = aspect3_df_text_diff)
    report.interviewee_gap_slide(template_slides['受訪者差異分析: 數位科技'], pic_rows=aspect4_qualitative_top5_gap, text_rows = aspect4_df_text_diff)
    
    """
    report.industry_slides(presentation, template_slides['客戶產業數位轉型重點與建議'], CASES_PER_SLIDE, industries, cases)    report.fin_indicator_slide(presentation, template_slides['財務指標表現'], INDICATOR_PER_SLIDE, rows=fin_indicator_data)
    report.fin_indicator_sensitivity_slide(template_slides['財務敏感度影響分析'], fin_sensitivity_plot_png)
    report.solution_priority_matrix_slide(template_slides['解決方案優先順序矩陣圖'], solution_priority_matrix_png)
    report.qualitative_questions_detail_slide(template_slides['質化題目填寫明細'], qualitative_plot_png)
    report.qualitative_plots_slide(presentation, template_slides['qualitative plots slide'], aspect_count, img2_buffers)
    """
    
    ppt_buffer = BytesIO()
    presentation.save(ppt_buffer)
    ppt_buffer.seek(0)
    
    #presentation.save('unused\\result.pptx')
    
    return ppt_buffer

    #presentation.save('test-output\\result.pptx')

    #img = Image.open(solution_priority_matrix_png)
    #img.show()
    #img.close()

    #print(df_trend)
    #print(df_fin_performance)
    #print(df_fin_sensitivity)
    #print(df_qualitative_result)
    #print(df_solution)

    #img = Image.open(qualitative_plot_png)
    #img.save('test-output\\qualitative.png')
    #img = Image.open(solution_priority_matrix_png)
    #img.save('test-output\\matrix.png')
    #img = Image.open(fin_sensitivity_plot_png)
    #img.save('test-output\\sensitivity.png')
    #img.close()
    # finally save file

    #img_buffer = fin_sensitivity_plot(df_fin_sensitivity)
